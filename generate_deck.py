import sys
import os

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("ERROR: python-pptx is not installed. Please run: pip install python-pptx")
    sys.exit(1)

def create_presentation():
    prs = Presentation()
    
    # Set to 16:9 widescreen format
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Styling Constants
    COLOR_BG_DEEP = RGBColor(7, 8, 20)          # #070814
    COLOR_BG_CARD = RGBColor(14, 16, 40)        # #0e1026
    COLOR_WHITE = RGBColor(248, 250, 252)       # #f8fafc
    COLOR_MUTED = RGBColor(148, 163, 184)       # #94a3b8
    COLOR_CYAN = RGBColor(6, 182, 212)          # #06b6d4
    COLOR_PURPLE = RGBColor(168, 85, 247)       # #a855f7
    COLOR_INDIGO = RGBColor(99, 102, 241)       # #6366f1
    COLOR_GREEN = RGBColor(16, 185, 129)        # #10b981
    COLOR_GOLD = RGBColor(251, 191, 36)         # #fbbf24
    
    FONT_TITLE = "Trebuchet MS"
    FONT_BODY = "Arial"
    
    # Helper: Set Slide Background
    def set_dark_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_BG_DEEP

    # Helper: Add standard Title text block
    def add_slide_header(slide, title_text, category_text="CORE MODULE"):
        # Category Tag
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(10), Inches(0.4))
        cat_tf = cat_box.text_frame
        cat_tf.word_wrap = True
        cat_tf.margin_left = Inches(0)
        cat_tf.margin_top = Inches(0)
        p_cat = cat_tf.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.name = FONT_TITLE
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = COLOR_CYAN
        
        # Main Header Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.5), Inches(0.8))
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        title_tf.margin_left = Inches(0)
        title_tf.margin_top = Inches(0)
        p_title = title_tf.paragraphs[0]
        p_title.text = title_text
        p_title.font.name = FONT_TITLE
        p_title.font.size = Pt(36)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_WHITE

    # Helper: Add premium card container
    def add_glass_card(slide, left, top, width, height, border_color=COLOR_INDIGO):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_BG_CARD
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
        return shape

    # ----------------------------------------------------
    # SLIDE 1: TITLE SLIDE (Premium & Glowy)
    # ----------------------------------------------------
    blank_slide_layout = prs.slide_layouts[6]
    slide1 = prs.slides.add_slide(blank_slide_layout)
    set_dark_background(slide1)
    
    # Glow Graphic (Mocked via shape)
    glow = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.6), Inches(1.5), Inches(6), Inches(4.5))
    glow.fill.solid()
    glow.fill.fore_color.rgb = COLOR_INDIGO
    glow.line.fill.background() # No border
    
    # Add a title container card to create depth
    add_glass_card(slide1, Inches(1.5), Inches(1.2), Inches(10.33), Inches(5.1), COLOR_CYAN)
    
    # Title Text Box
    t_box = slide1.shapes.add_textbox(Inches(2.0), Inches(1.8), Inches(9.33), Inches(3.8))
    tf = t_box.text_frame
    tf.word_wrap = True
    
    p0 = tf.paragraphs[0]
    p0.text = "SKILLFORGE AI"
    p0.font.name = FONT_TITLE
    p0.font.size = Pt(64)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_WHITE
    p0.alignment = PP_ALIGN.CENTER
    
    p1 = tf.add_paragraph()
    p1.text = "Transforming Degrees Into Active Careers"
    p1.font.name = FONT_TITLE
    p1.font.size = Pt(28)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_CYAN
    p1.alignment = PP_ALIGN.CENTER
    p1.space_before = Pt(14)
    
    p2 = tf.add_paragraph()
    p2.text = "An AI-Powered Employability Diagnostics, Adaptable Roadmaps, ATS Resume Optimizers,\nMock Technical Interviews, and Entrepreneur Launchpad."
    p2.font.name = FONT_BODY
    p2.font.size = Pt(14)
    p2.font.color.rgb = COLOR_MUTED
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(24)
    
    p3 = tf.add_paragraph()
    p3.text = "NATIONAL HACKATHON SHOWCASE  •  TEAM INCUBATOR"
    p3.font.name = FONT_TITLE
    p3.font.size = Pt(11)
    p3.font.bold = True
    p3.font.color.rgb = COLOR_PURPLE
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(40)

    # ----------------------------------------------------
    # SLIDE 2: THE PROBLEM (Academic & Industry Chasm)
    # ----------------------------------------------------
    slide2 = prs.slides.add_slide(blank_slide_layout)
    set_dark_background(slide2)
    add_slide_header(slide2, "The Employability Chasm", "THE CORE PROBLEM")
    
    # 3 Cards for Problems
    card_w = Inches(3.6)
    card_h = Inches(4.3)
    card_y = Inches(2.0)
    
    # Problem 1
    add_glass_card(slide2, Inches(0.8), card_y, card_w, card_h, COLOR_PURPLE)
    p1_box = slide2.shapes.add_textbox(Inches(1.0), card_y + Inches(0.2), card_w - Inches(0.4), card_h - Inches(0.4))
    tf1 = p1_box.text_frame
    tf1.word_wrap = True
    
    p_h1 = tf1.paragraphs[0]
    p_h1.text = "01. Academic Inertia"
    p_h1.font.name = FONT_TITLE
    p_h1.font.size = Pt(22)
    p_h1.font.bold = True
    p_h1.font.color.rgb = COLOR_PURPLE
    
    p_b1 = tf1.add_paragraph()
    p_b1.text = "\nAcademic syllabi adapt extremely slowly to fast-paced technical advancements. Students learn standard theories but miss production patterns like REST routing, containerization, and modern state frameworks."
    p_b1.font.name = FONT_BODY
    p_b1.font.size = Pt(13)
    p_b1.font.color.rgb = COLOR_WHITE
    p_b1.space_before = Pt(10)
    
    # Problem 2
    add_glass_card(slide2, Inches(4.8), card_y, card_w, card_h, COLOR_PURPLE)
    p2_box = slide2.shapes.add_textbox(Inches(5.0), card_y + Inches(0.2), card_w - Inches(0.4), card_h - Inches(0.4))
    tf2 = p2_box.text_frame
    tf2.word_wrap = True
    
    p_h2 = tf2.paragraphs[0]
    p_h2.text = "02. The Certification Trap"
    p_h2.font.name = FONT_TITLE
    p_h2.font.size = Pt(22)
    p_h2.font.bold = True
    p_h2.font.color.rgb = COLOR_PURPLE
    
    p_b2 = tf2.add_paragraph()
    p_b2.text = "\nStudents collect generic video course certifications without ever writing custom files or debugging server faults. When asked to construct real codebases in technical checks, they fail to deliver."
    p_b2.font.name = FONT_BODY
    p_b2.font.size = Pt(13)
    p_b2.font.color.rgb = COLOR_WHITE
    p_b2.space_before = Pt(10)
    
    # Problem 3
    add_glass_card(slide2, Inches(8.8), card_y, card_w, card_h, COLOR_PURPLE)
    p3_box = slide2.shapes.add_textbox(Inches(9.0), card_y + Inches(0.2), card_w - Inches(0.4), card_h - Inches(0.4))
    tf3 = p3_box.text_frame
    tf3.word_wrap = True
    
    p_h3 = tf3.paragraphs[0]
    p_h3.text = "03. Institutional Blindness"
    p_h3.font.name = FONT_TITLE
    p_h3.font.size = Pt(22)
    p_h3.font.bold = True
    p_h3.font.color.rgb = COLOR_PURPLE
    
    p_b3 = tf3.add_paragraph()
    p_b3.text = "\nPlacement officers and college deans have zero live metrics on department capability vectors. They coordinate campus placement drives blindly, resulting in poor conversion rates and recruiter dissatisfaction."
    p_b3.font.name = FONT_BODY
    p_b3.font.size = Pt(13)
    p_b3.font.color.rgb = COLOR_WHITE
    p_b3.space_before = Pt(10)

    # ----------------------------------------------------
    # SLIDE 3: THE SOLUTION (SkillForge AI Platform)
    # ----------------------------------------------------
    slide3 = prs.slides.add_slide(blank_slide_layout)
    set_dark_background(slide3)
    add_slide_header(slide3, "The Active Solution: SkillForge AI", "THE RESOLUTION")
    
    card_w = Inches(5.6)
    card_h = Inches(2.0)
    
    # Solution 1
    add_glass_card(slide3, Inches(0.8), Inches(2.0), card_w, card_h, COLOR_CYAN)
    s1_box = slide3.shapes.add_textbox(Inches(1.0), Inches(2.1), card_w - Inches(0.4), card_h - Inches(0.2))
    s1_tf = s1_box.text_frame
    s1_tf.word_wrap = True
    p_sh1 = s1_tf.paragraphs[0]
    p_sh1.text = "AI Skill Gap Diagnostics"
    p_sh1.font.name = FONT_TITLE
    p_sh1.font.size = Pt(18)
    p_sh1.font.bold = True
    p_sh1.font.color.rgb = COLOR_CYAN
    p_sb1 = s1_tf.add_paragraph()
    p_sb1.text = "Precision audits comparing custom job descriptions against the student's background using real Gemini AI. Adaptively updates timeline curricula instantly."
    p_sb1.font.name = FONT_BODY
    p_sb1.font.size = Pt(12)
    p_sb1.font.color.rgb = COLOR_MUTED
    
    # Solution 2
    add_glass_card(slide3, Inches(6.8), Inches(2.0), card_w, card_h, COLOR_CYAN)
    s2_box = slide3.shapes.add_textbox(Inches(7.0), Inches(2.1), card_w - Inches(0.4), card_h - Inches(0.2))
    s2_tf = s2_box.text_frame
    s2_tf.word_wrap = True
    p_sh2 = s2_tf.paragraphs[0]
    p_sh2.text = "Dynamic Milestones & Projects"
    p_sh2.font.name = FONT_TITLE
    p_sh2.font.size = Pt(18)
    p_sh2.font.bold = True
    p_sh2.font.color.rgb = COLOR_CYAN
    p_sb2 = s2_tf.add_paragraph()
    p_sb2.text = "Week-by-week interactive timeline milestones tied to core sandbox projects. Instantly recalculates capability indexes on task completions."
    p_sb2.font.name = FONT_BODY
    p_sb2.font.size = Pt(12)
    p_sb2.font.color.rgb = COLOR_MUTED
    
    # Solution 3
    add_glass_card(slide3, Inches(0.8), Inches(4.5), card_w, card_h, COLOR_CYAN)
    s3_box = slide3.shapes.add_textbox(Inches(1.0), Inches(4.6), card_w - Inches(0.4), card_h - Inches(0.2))
    s3_tf = s3_box.text_frame
    s3_tf.word_wrap = True
    p_sh3 = s3_tf.paragraphs[0]
    p_sh3.text = "ATS Scanner & Mock Assessment"
    p_sh3.font.name = FONT_TITLE
    p_sh3.font.size = Pt(18)
    p_sh3.font.bold = True
    p_sh3.font.color.rgb = COLOR_CYAN
    p_sb3 = s3_tf.add_paragraph()
    p_sb3.text = "Live lexical scanners that evaluate ATS keyword match indices and offer interactive AI interview rounds with scoring feedbacks."
    p_sb3.font.name = FONT_BODY
    p_sb3.font.size = Pt(12)
    p_sb3.font.color.rgb = COLOR_MUTED
    
    # Solution 4
    add_glass_card(slide3, Inches(6.8), Inches(4.5), card_w, card_h, COLOR_CYAN)
    s4_box = slide3.shapes.add_textbox(Inches(7.0), Inches(4.6), card_w - Inches(0.4), card_h - Inches(0.2))
    s4_tf = s4_box.text_frame
    s4_tf.word_wrap = True
    p_sh4 = s4_tf.paragraphs[0]
    p_sh4.text = "Startup Studio & Incubation"
    p_sh4.font.name = FONT_TITLE
    p_sh4.font.size = Pt(18)
    p_sh4.font.bold = True
    p_sh4.font.color.rgb = COLOR_CYAN
    p_sb4 = s4_tf.add_paragraph()
    p_sb4.text = "Accelerates startup concepts, instantly generating complete 9-box Business Model Canvas frameworks details via Gemini integrations."
    p_sb4.font.name = FONT_BODY
    p_sb4.font.size = Pt(12)
    p_sb4.font.color.rgb = COLOR_MUTED

    # ----------------------------------------------------
    # SLIDE 4: THE TECH STACK & LOCAL STORAGE DATABASE
    # ----------------------------------------------------
    slide4 = prs.slides.add_slide(blank_slide_layout)
    set_dark_background(slide4)
    add_slide_header(slide4, "Under the Hood: High-Fidelity Local Tech Stack", "SYSTEM ARCHITECTURE")
    
    # Tech Stack Box Left
    add_glass_card(slide4, Inches(0.8), Inches(2.0), Inches(5.6), Inches(4.5), COLOR_INDIGO)
    stack_box = slide4.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(5.2), Inches(4.1))
    st_tf = stack_box.text_frame
    st_tf.word_wrap = True
    p_sth = st_tf.paragraphs[0]
    p_sth.text = "Engineered Architecture"
    p_sth.font.name = FONT_TITLE
    p_sth.font.size = Pt(22)
    p_sth.font.bold = True
    p_sth.font.color.rgb = COLOR_INDIGO
    
    p_stb = st_tf.add_paragraph()
    p_stb.text = "\n•  UI Layer: React 18 SPA built with Vite for extreme sub-second hot modules reloading.\n\n•  Styling: Pure, customized Vanilla CSS utilizing variables. Deep HSL palette, dark space backdrop with glassmorphic cards and dynamic flex panels.\n\n•  Visuals: Recharts analytics for responsive, animated capability vectors (Bar and Polar Radar models).\n\n•  AI Core: Real REST endpoints structured to query Gemini 2.0 Flash for low-latency JSON response templates."
    p_stb.font.name = FONT_BODY
    p_stb.font.size = Pt(13)
    p_stb.font.color.rgb = COLOR_WHITE
    p_stb.space_before = Pt(8)
    
    # Database Box Right
    add_glass_card(slide4, Inches(6.8), Inches(2.0), Inches(5.6), Inches(4.5), COLOR_INDIGO)
    db_box = slide4.shapes.add_textbox(Inches(7.0), Inches(2.2), Inches(5.2), Inches(4.1))
    db_tf = db_box.text_frame
    db_tf.word_wrap = True
    p_dbh = db_tf.paragraphs[0]
    p_dbh.text = "Persistent Database Engine"
    p_dbh.font.name = FONT_TITLE
    p_dbh.font.size = Pt(22)
    p_dbh.font.bold = True
    p_dbh.font.color.rgb = COLOR_INDIGO
    
    p_dbb = db_tf.add_paragraph()
    p_dbb.text = "\n•  Simulated Multi-User Database: We bypassed standard heavy cloud solutions to optimize live presentation reliability.\n\n•  Localized JSON Store: Created a custom Repository Pattern (`dbService.js`) syncing student contexts to `localStorage`.\n\n•  Automatic Hydration: Tracks account credentials, custom profile parameters, checked roadmap items, badge statuses, and generated canvas data across sessions.\n\n•  Enterprise Swapping: Binds frontend contexts solely to db interfaces, enabling painless migration to PostgreSQL/Node.js."
    p_dbb.font.name = FONT_BODY
    p_dbb.font.size = Pt(13)
    p_dbb.font.color.rgb = COLOR_WHITE
    p_dbb.space_before = Pt(8)

    # ----------------------------------------------------
    # SLIDE 5: PILLAR 1 - AI SKILL GAP ANALYZER
    # ----------------------------------------------------
    slide5 = prs.slides.add_slide(blank_slide_layout)
    set_dark_background(slide5)
    add_slide_header(slide5, "Pillar 01: Adaptive AI Skill Gap Analyzer", "MODULE DEMO")
    
    # Mock Graphic Left (SVG Dial)
    add_glass_card(slide5, Inches(0.8), Inches(2.0), Inches(4.8), Inches(4.5), COLOR_CYAN)
    g_box = slide5.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(4.4), Inches(4.1))
    g_tf = g_box.text_frame
    g_tf.word_wrap = True
    p_gh = g_tf.paragraphs[0]
    p_gh.text = "Real-Time AI Audits"
    p_gh.font.name = FONT_TITLE
    p_gh.font.size = Pt(20)
    p_gh.font.bold = True
    p_gh.font.color.rgb = COLOR_CYAN
    p_gb = g_tf.add_paragraph()
    p_gb.text = "\n1. Inputs: Target career path + existing credentials.\n\n2. Engine: Structural prompt calls Gemini REST API.\n\n3. Output: Generates missing vs existing arrays.\n\n4. Integration: Adaptively updates the curriculum timeline checklists in real time!"
    p_gb.font.name = FONT_BODY
    p_gb.font.size = Pt(14)
    p_gb.font.color.rgb = COLOR_WHITE
    
    # Detailed Description Right
    desc_box = slide5.shapes.add_textbox(Inches(6.0), Inches(2.0), Inches(6.5), Inches(4.5))
    d_tf = desc_box.text_frame
    d_tf.word_wrap = True
    p_dh = d_tf.paragraphs[0]
    p_dh.text = "How the Engine Maps the Gaps"
    p_dh.font.name = FONT_TITLE
    p_dh.font.size = Pt(24)
    p_dh.font.bold = True
    p_dh.font.color.rgb = COLOR_WHITE
    
    p_db = d_tf.add_paragraph()
    p_db.text = "SkillForge AI solves the standard certification mismatch. Instead of forcing students through generic pathways, our Gemini analyzer audits their target role."
    p_db.font.name = FONT_BODY
    p_db.font.size = Pt(14)
    p_db.font.color.rgb = COLOR_MUTED
    p_db.space_before = Pt(8)
    
    bullet_box = d_tf.add_paragraph()
    bullet_box.text = "•  Structured JSON Output: The AI formats outputs with structured keys, allowing the React framework to map values immediately without regex breaks.\n•  Targeted Focus Board: Isolates missing core skills (e.g. JWT Auth, Relational joins) so students avoid repeating basics.\n•  Curriculum Injector: Direct programmatic loops feed the extracted missing parameters directly into Roadmap Weeks 4-6."
    bullet_box.font.name = FONT_BODY
    bullet_box.font.size = Pt(13)
    bullet_box.font.color.rgb = COLOR_WHITE
    bullet_box.space_before = Pt(12)

    # ----------------------------------------------------
    # SLIDE 6: PILLAR 2 - ROADMAPS & SANDBOX PROJECTS
    # ----------------------------------------------------
    slide6 = prs.slides.add_slide(blank_slide_layout)
    set_dark_background(slide6)
    add_slide_header(slide6, "Pillar 02: Adaptive Roadmaps & Project Boards", "MODULE DEMO")
    
    # Description Left
    r_desc_box = slide6.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.6), Inches(4.5))
    rd_tf = r_desc_box.text_frame
    rd_tf.word_wrap = True
    p_rdh = rd_tf.paragraphs[0]
    p_rdh.text = "Practical Milestone Timelines"
    p_rdh.font.name = FONT_TITLE
    p_rdh.font.size = Pt(24)
    p_rdh.font.bold = True
    p_rdh.font.color.rgb = COLOR_WHITE
    
    p_rdb = rd_tf.add_paragraph()
    p_rdb.text = "Theoretical learning fails recruiters. SkillForge enforces progress through interactive week-by-week timelines linked directly to sandbox coding challenges."
    p_rdb.font.name = FONT_BODY
    p_rdb.font.size = Pt(14)
    p_rdb.font.color.rgb = COLOR_MUTED
    p_rdb.space_before = Pt(8)
    
    p_rdbul = rd_tf.add_paragraph()
    p_rdbul.text = "•  Locked Milestones: Weekly syllabi auto-unlock as earlier checks are completed, driving progressive achievement.\n•  XP Scoring Loops: Completed roadmap checks immediately update the student's Project and Readiness scores in local database arrays.\n•  Verification Gates: Certificate unlocking acts as an incentive, forcing students to hit the critical 80% capability index."
    p_rdbul.font.name = FONT_BODY
    p_rdbul.font.size = Pt(13)
    p_rdbul.font.color.rgb = COLOR_WHITE
    p_rdbul.space_before = Pt(12)
    
    # Cards Right (Mock Project card)
    add_glass_card(slide6, Inches(6.8), Inches(2.0), Inches(5.6), Inches(4.5), COLOR_INDIGO)
    proj_box = slide6.shapes.add_textbox(Inches(7.0), Inches(2.2), Inches(5.2), Inches(4.1))
    pj_tf = proj_box.text_frame
    pj_tf.word_wrap = True
    p_pjh = pj_tf.paragraphs[0]
    p_pjh.text = "AI Project Sandbox Cards"
    p_pjh.font.name = FONT_TITLE
    p_pjh.font.size = Pt(20)
    p_pjh.font.bold = True
    p_pjh.font.color.rgb = COLOR_INDIGO
    
    p_pjb = pj_tf.add_paragraph()
    p_pjb.text = "\n•  Beginner: Custom weather scans and to-do workspaces utilizing local persistence layers.\n\n•  Intermediate: Relational expense journals and web-socket messenger systems using Express routers.\n\n•  Advanced: Automated NLP ATS scanners and real-time mock interview chatbot interfaces.\n\n•  Integrated Checkpoints: Completing coding sandboxes triggers a notification toast and updates overall employability indices."
    p_pjb.font.name = FONT_BODY
    p_pjb.font.size = Pt(13)
    p_pjb.font.color.rgb = COLOR_WHITE

    # ----------------------------------------------------
    # SLIDE 7: PILLAR 3 - ATS SCANNER & INTERVIEWS
    # ----------------------------------------------------
    slide7 = prs.slides.add_slide(blank_slide_layout)
    set_dark_background(slide7)
    add_slide_header(slide7, "Pillar 03: NLP Resume Scanners & Mock Interview", "MODULE DEMO")
    
    # Resume Box Left
    add_glass_card(slide7, Inches(0.8), Inches(2.0), Inches(5.6), Inches(4.5), COLOR_GREEN)
    res_box = slide7.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(5.2), Inches(4.1))
    rs_tf = res_box.text_frame
    rs_tf.word_wrap = True
    p_rsh = rs_tf.paragraphs[0]
    p_rsh.text = "Simulated ATS Scanners"
    p_rsh.font.name = FONT_TITLE
    p_rsh.font.size = Pt(22)
    p_rsh.font.bold = True
    p_rsh.font.color.rgb = COLOR_GREEN
    
    p_rsb = rs_tf.add_paragraph()
    p_rsb.text = "\n•  Keyword Parser: Running a simulated parser evaluating resume drafts against target career criteria.\n\n•  Actionable Weaknesses: Isolates missing tech terms (e.g. Git, Express) and highlights formatting issues.\n\n•  Auto-Injection Panel: Programmatic keywords injector enables students to sync resume vocabularies, instantly boosting ATS compatibility ratings and indices."
    p_rsb.font.name = FONT_BODY
    p_rsb.font.size = Pt(13)
    p_rsb.font.color.rgb = COLOR_WHITE
    p_rsb.space_before = Pt(8)
    
    # Interview Box Right
    add_glass_card(slide7, Inches(6.8), Inches(2.0), Inches(5.6), Inches(4.5), COLOR_GREEN)
    int_box = slide7.shapes.add_textbox(Inches(7.0), Inches(2.2), Inches(5.2), Inches(4.1))
    it_tf = int_box.text_frame
    it_tf.word_wrap = True
    p_ith = it_tf.paragraphs[0]
    p_ith.text = "Mock Interview Boards"
    p_ith.font.name = FONT_TITLE
    p_ith.font.size = Pt(22)
    p_ith.font.bold = True
    p_ith.font.color.rgb = COLOR_GREEN
    
    p_itb = it_tf.add_paragraph()
    p_itb.text = "\n•  Dynamic Loops: Students practice core Technical, HR, and behavioral STAR-format interview questions.\n\n•  Interactive Chatbot: Binds response inputs to the floating Gemini companion, generating highly custom counter-questions.\n\n•  Performance Metric: Delivers analytical performance scoring cards detailing answer strengths and grammatical recommendations."
    p_itb.font.name = FONT_BODY
    p_itb.font.size = Pt(13)
    p_itb.font.color.rgb = COLOR_WHITE
    p_itb.space_before = Pt(8)

    # ----------------------------------------------------
    # SLIDE 8: PILLAR 4 - ENTREPRENEUR LAUNCHPAD
    # ----------------------------------------------------
    slide8 = prs.slides.add_slide(blank_slide_layout)
    set_dark_background(slide8)
    add_slide_header(slide8, "Pillar 04: Entrepreneur Startup Launchpad", "MODULE DEMO")
    
    # Description Left
    ent_desc = slide8.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.6), Inches(4.5))
    en_tf = ent_desc.text_frame
    en_tf.word_wrap = True
    p_enh = en_tf.paragraphs[0]
    p_enh.text = "Incubating Campus Founders"
    p_enh.font.name = FONT_TITLE
    p_enh.font.size = Pt(24)
    p_enh.font.bold = True
    p_enh.font.color.rgb = COLOR_WHITE
    
    p_enb = en_tf.add_paragraph()
    p_enb.text = "Placement is only half the battle. Our Entrepreneur Launchpad caters to student founders, helping them structure business models and identify market fits on day one."
    p_enb.font.name = FONT_BODY
    p_enb.font.size = Pt(14)
    p_enb.font.color.rgb = COLOR_MUTED
    p_enb.space_before = Pt(8)
    
    p_enbul = en_tf.add_paragraph()
    p_enbul.text = "•  AI Model Canvas: Generates 9 distinct fields outlining partners, customer segments, channels, and costs.\n•  Structured Pitch Decks: Outlines exact slide frameworks and key talking points to present to incubator grant panels.\n•  Incubation XP: Fostering entrepreneurial tasks scales the global portfolio score, validating versatile engineering talents."
    p_enbul.font.name = FONT_BODY
    p_enbul.font.size = Pt(13)
    p_enbul.font.color.rgb = COLOR_WHITE
    p_enbul.space_before = Pt(12)
    
    # Business Canvas Visual Right
    add_glass_card(slide8, Inches(6.8), Inches(2.0), Inches(5.6), Inches(4.5), COLOR_GOLD)
    canvas_box = slide8.shapes.add_textbox(Inches(7.0), Inches(2.2), Inches(5.2), Inches(4.1))
    cv_tf = canvas_box.text_frame
    cv_tf.word_wrap = True
    p_cvh = cv_tf.paragraphs[0]
    p_cvh.text = "AI Business Canvas Layout"
    p_cvh.font.name = FONT_TITLE
    p_cvh.font.size = Pt(20)
    p_cvh.font.bold = True
    p_cvh.font.color.rgb = COLOR_GOLD
    
    p_cvb = cv_tf.add_paragraph()
    p_cvb.text = "\n•  Value Propositions: Details unique customer problem-solving solutions generated in real time.\n\n•  Key Channels: Outlines targeted advertising campaigns and distribution systems.\n\n•  Cost Structures: Maps initial cloud services, hosting, and human resource parameters.\n\n•  Revenue Streams: Designs structured SaaS, premium licensing, or transaction fees models."
    p_cvb.font.name = FONT_BODY
    p_cvb.font.size = Pt(13)
    p_cvb.font.color.rgb = COLOR_WHITE

    # ----------------------------------------------------
    # SLIDE 9: COLLEGE PORTAL & ADMIN TELEMETRY
    # ----------------------------------------------------
    slide9 = prs.slides.add_slide(blank_slide_layout)
    set_dark_background(slide9)
    add_slide_header(slide9, "Campus College & Platform Admin Dashboards", "ENTERPRISE B2B SYSTEMS")
    
    # B2B Left
    add_glass_card(slide9, Inches(0.8), Inches(2.0), Inches(5.6), Inches(4.5), COLOR_INDIGO)
    col_box = slide9.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(5.2), Inches(4.1))
    cl_tf = col_box.text_frame
    cl_tf.word_wrap = True
    p_clh = cl_tf.paragraphs[0]
    p_clh.text = "Collegiate Placements Portal"
    p_clh.font.name = FONT_TITLE
    p_clh.font.size = Pt(22)
    p_clh.font.bold = True
    p_clh.font.color.rgb = COLOR_INDIGO
    
    p_clb = cl_tf.add_paragraph()
    p_clb.text = "\n•  Aggregated Metrics: College Placement Officers scan department-wide capability matrices (CSE vs IT).\n\n•  Job-Ready Filters: Instantly identifies high-scoring students (80%+) to expedite recruiter invitation batches.\n\n•  Skill Trend Analytics: Visualizes department readiness tracks, helping universities adapt academic resources."
    p_clb.font.name = FONT_BODY
    p_clb.font.size = Pt(13)
    p_clb.font.color.rgb = COLOR_WHITE
    p_clb.space_before = Pt(8)
    
    # Admin Right
    add_glass_card(slide9, Inches(6.8), Inches(2.0), Inches(5.6), Inches(4.5), COLOR_INDIGO)
    adm_box = slide9.shapes.add_textbox(Inches(7.0), Inches(2.2), Inches(5.2), Inches(4.1))
    ad_tf = adm_box.text_frame
    ad_tf.word_wrap = True
    p_adh = ad_tf.paragraphs[0]
    p_adh.text = "Platform Security & Telemetry"
    p_adh.font.name = FONT_TITLE
    p_adh.font.size = Pt(22)
    p_adh.font.bold = True
    p_adh.font.color.rgb = COLOR_INDIGO
    
    p_adb = ad_tf.add_paragraph()
    p_adb.text = "\n•  Diagnostics Log Box: Simulates low-level platform server telemetry in real-time script templates.\n\n•  Key Stats: Monitors total student registrations, API latency values, and active storage allocations.\n\n•  Perspective Toggle: Allows judges to dynamically swap perspectives (Student, Dean, Admin) to experience three platforms in one."
    p_adb.font.name = FONT_BODY
    p_adb.font.size = Pt(13)
    p_adb.font.color.rgb = COLOR_WHITE
    p_adb.space_before = Pt(8)

    # ----------------------------------------------------
    # SLIDE 10: WHY WE WIN (THE HACKATHON WINNER'S EDGE)
    # ----------------------------------------------------
    slide10 = prs.slides.add_slide(blank_slide_layout)
    set_dark_background(slide10)
    add_slide_header(slide10, "Why SkillForge AI Wins This Hackathon", "THE WINNING FORMULA")
    
    # Left Card
    add_glass_card(slide10, Inches(0.8), Inches(2.0), Inches(5.6), Inches(4.5), COLOR_GOLD)
    win_box1 = slide10.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(5.2), Inches(4.1))
    w1_tf = win_box1.text_frame
    w1_tf.word_wrap = True
    p_w1h = w1_tf.paragraphs[0]
    p_w1h.text = "Judge-Centric Integrations"
    p_w1h.font.name = FONT_TITLE
    p_w1h.font.size = Pt(22)
    p_w1h.font.bold = True
    p_w1h.font.color.rgb = COLOR_GOLD
    
    p_w1b = w1_tf.add_paragraph()
    p_w1b.text = "\n•  Pitch Mode Superimposition: Programmatic overlay cards highlight design rationale, technical architecture choices, and presentation tips.\n\n•  Zero Layout Placeholders: Over 10,000 words of rich, realistic context. All chats, canvases, roadmaps, and diagnostics display fully detailed fields.\n\n•  Three Platforms in One: Features Student, College Dean, and System Admin perspectives in a single deployment."
    p_w1b.font.name = FONT_BODY
    p_w1b.font.size = Pt(13)
    p_w1b.font.color.rgb = COLOR_WHITE
    p_w1b.space_before = Pt(8)
    
    # Right Card
    add_glass_card(slide10, Inches(6.8), Inches(2.0), Inches(5.6), Inches(4.5), COLOR_GOLD)
    win_box2 = slide10.shapes.add_textbox(Inches(7.0), Inches(2.2), Inches(5.2), Inches(4.1))
    w2_tf = win_box2.text_frame
    w2_tf.word_wrap = True
    p_w2h = w2_tf.paragraphs[0]
    p_w2h.text = "Technical & Product Viability"
    p_w2h.font.name = FONT_TITLE
    p_w2h.font.size = Pt(22)
    p_w2h.font.bold = True
    p_w2h.font.color.rgb = COLOR_GOLD
    
    p_w2b = w2_tf.add_paragraph()
    p_w2b.text = "\n•  State Reactivity: Binds progress checklist updates to capability vector meters via custom calculations.\n\n•  Scalability Blueprint: Decoupled design patterns make frontend modules entirely independent of storage models.\n\n•  Real AI Orchestration: Real-time REST queries connect to Gemini 2.0 Flash, delivering fast results and complete business models."
    p_w2b.font.name = FONT_BODY
    p_w2b.font.size = Pt(13)
    p_w2b.font.color.rgb = COLOR_WHITE
    p_w2b.space_before = Pt(8)
    
    # Save Presentation
    prs.save("SkillForge_AI_Hackathon_Pitch.pptx")
    print("SUCCESS: SkillForge_AI_Hackathon_Pitch.pptx generated successfully in widescreen layout!")

if __name__ == "__main__":
    create_presentation()
